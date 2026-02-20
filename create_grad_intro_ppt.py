from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Color Palette: Midnight Executive
COLOR_NAVY = RGBColor(30, 39, 97)      # #1E2761
COLOR_ICE_BLUE = RGBColor(202, 220, 252) # #CADCFC
COLOR_WHITE = RGBColor(255, 255, 255)    # #FFFFFF
COLOR_CHARCOAL = RGBColor(54, 69, 79)    # #36454F
COLOR_ACCENT = RGBColor(255, 215, 0)     # #FFD700 (Gold for accents)

def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, COLOR_NAVY)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    tf.text = subtitle_text
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_ICE_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Decorative line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(4.2), Inches(4), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

def add_content_slide(prs, title_text, content_list):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, COLOR_WHITE)

    # Title Bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_NAVY
    shape.line.fill.background()

    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_CHARCOAL
        p.space_after = Pt(14)
        p.level = 0

def add_two_column_slide(prs, title_text, left_content, right_content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, COLOR_WHITE)

    # Title Bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_NAVY
    shape.line.fill.background()

    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Left Column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for item in left_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_CHARCOAL
        p.space_after = Pt(10)

    # Right Column
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for item in right_content:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_CHARCOAL
        p.space_after = Pt(10)

def add_thank_you_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, COLOR_NAVY)

    # Main Text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "感谢聆听 敬请指正"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Contact Info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tf = contact_box.text_frame
    tf.text = "汇报人：[您的姓名]\n邮箱：[您的邮箱]"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_ICE_BLUE
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Cover
    add_title_slide(prs, "研究生入学自我介绍", "汇报人：[您的姓名]\n导师：[导师姓名]")

    # Slide 2: Introduction (Two Columns)
    add_two_column_slide(prs, "个人简介", 
        ["姓名：[您的姓名]", "本科学校：[您的本科学校]", "本科专业：[您的专业]", "GPA：3.8/4.0 (前 5%)", "家乡：[省份城市]"],
        ["此处可放置个人照片", " ", "座右铭：", "“脚踏实地，仰望星空”", " ", "兴趣爱好：", "阅读、编程、羽毛球"]
    )

    # Slide 3: Academic Experience
    add_content_slide(prs, "学术/科研经历", [
        "• 项目经历一：[项目名称]",
        "  - 担任角色：核心成员",
        "  - 主要工作：负责数据预处理与模型搭建",
        "  - 成果：发表会议论文一篇",
        " ",
        "• 项目经历二：[项目名称]",
        "  - 担任角色：项目负责人",
        "  - 主要工作：统筹项目进度，撰写结题报告",
        "  - 成果：获得校级优秀项目奖"
    ])

    # Slide 4: Honors & Awards
    add_content_slide(prs, "荣誉奖项", [
        "• 202X年 国家奖学金",
        "• 202X年 全国大学生数学建模竞赛 一等奖",
        "• 202X年 “挑战杯”大学生课外学术科技作品竞赛 省级二等奖",
        "• 202X年 校级三好学生标兵",
        "• 202X年 英语六级 (CET-6): 600分"
    ])

    # Slide 5: Professional Skills
    add_two_column_slide(prs, "专业技能", 
        ["编程语言：", "• Python (熟练)", "• C++ (良好)", "• Java (基础)", " ", "开发工具：", "• PyCharm, VS Code, Git"],
        ["科研工具：", "• LaTeX (熟练)", "• MATLAB (良好)", "• SPSS (基础)", " ", "深度学习框架：", "• PyTorch", "• TensorFlow"]
    )

    # Slide 6: Research Plan
    add_content_slide(prs, "研究生规划", [
        "• 研一阶段：打好基础",
        "  - 完成研究生课程学习，保持优异成绩",
        "  - 广泛阅读领域文献，确定具体研究方向",
        " ",
        "• 研二阶段：深入科研",
        "  - 开展课题研究，进行实验验证",
        "  - 撰写并发表高水平学术论文",
        " ",
        "• 研三阶段：总结与展望",
        "  - 完成毕业论文撰写与答辩",
        "  - 规划职业发展，准备深造或就业"
    ])

    # Slide 7: Thank You
    add_thank_you_slide(prs)

    output_file = "Graduate_Introduction.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation()
